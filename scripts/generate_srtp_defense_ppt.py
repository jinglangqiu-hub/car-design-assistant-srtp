from __future__ import annotations

import csv
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "重庆大学SRTP结项答辩PPT_汽车外观设计辅助系统.pptx"

CQU_RED = RGBColor(151, 31, 45)
DEEP_BLUE = RGBColor(22, 45, 78)
INK = RGBColor(31, 41, 55)
MUTED = RGBColor(100, 116, 139)
PAPER = RGBColor(248, 250, 252)
PALE = RGBColor(239, 243, 248)
GOLD = RGBColor(196, 144, 63)
WHITE = RGBColor(255, 255, 255)


def add_shape(slide, shape_type, x, y, w, h, fill, line=None, transparency=0):
    shp = slide.shapes.add_shape(shape_type, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.fill.transparency = transparency
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    return shp


def add_text(slide, text, x, y, w, h, size=20, bold=False, color=INK, font="Microsoft YaHei UI",
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_multiline(slide, lines, x, y, w, h, size=17, color=INK, bullet=False, gap=4):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(gap)
        if bullet:
            p.text = f"• {line}"
        else:
            p.text = line
        for run in p.runs:
            run.font.name = "Microsoft YaHei UI"
            run.font.size = Pt(size)
            run.font.color.rgb = color
    return box


def set_bg(slide, color=PAPER):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_header(slide, no, title, kicker="SRTP 结项答辩"):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.14), CQU_RED)
    add_text(slide, kicker, Inches(0.55), Inches(0.34), Inches(2.4), Inches(0.25), 9, True, CQU_RED)
    add_text(slide, title, Inches(0.55), Inches(0.58), Inches(10.1), Inches(0.55), 25, True, DEEP_BLUE)
    add_text(slide, f"{no:02d}", Inches(11.85), Inches(0.42), Inches(0.7), Inches(0.35), 14, True, CQU_RED, align=PP_ALIGN.RIGHT)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(1.15), Inches(11.95), Inches(0.015), GOLD)


def add_footer(slide):
    add_text(slide, "重庆大学 SRTP | 汽车外观设计辅助系统", Inches(0.55), Inches(7.08), Inches(4.0), Inches(0.18), 8, False, MUTED)


def add_card(slide, x, y, w, h, title, body, accent=CQU_RED):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h, WHITE, RGBColor(220, 226, 235))
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, Inches(0.07), h, accent)
    add_text(slide, title, x + Inches(0.22), y + Inches(0.15), w - Inches(0.35), Inches(0.28), 15, True, DEEP_BLUE)
    add_multiline(slide, body, x + Inches(0.22), y + Inches(0.55), w - Inches(0.35), h - Inches(0.7), 12.5, INK)


def add_metric(slide, x, y, w, h, value, label, note="", color=CQU_RED):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h, WHITE, RGBColor(222, 226, 232))
    add_text(slide, value, x + Inches(0.12), y + Inches(0.16), w - Inches(0.24), Inches(0.38), 23, True, color, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + Inches(0.12), y + Inches(0.62), w - Inches(0.24), Inches(0.25), 10.5, True, INK, align=PP_ALIGN.CENTER)
    if note:
        add_text(slide, note, x + Inches(0.12), y + Inches(0.9), w - Inches(0.24), Inches(0.25), 8.5, False, MUTED, align=PP_ALIGN.CENTER)


def add_arrow(slide, x1, y1, x2, y2, color=MUTED):
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x1, y1 - Inches(0.01), x2 - x1 - Inches(0.08), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    tri = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_TRIANGLE, x2 - Inches(0.11), y2 - Inches(0.07), Inches(0.13), Inches(0.14))
    tri.fill.solid()
    tri.fill.fore_color.rgb = color
    tri.line.fill.background()
    return line


def add_picture_fit(slide, img_path, x, y, w, h):
    if not img_path.exists():
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h, PALE, RGBColor(220, 226, 235))
        add_text(slide, "图片待补充", x, y + h / 2 - Inches(0.15), w, Inches(0.3), 14, True, MUTED, align=PP_ALIGN.CENTER)
        return
    pic = slide.shapes.add_picture(str(img_path), x, y, width=w)
    if pic.height > h:
        scale = h / pic.height
        pic.width = int(pic.width * scale)
        pic.height = h
    pic.left = int(x + (w - pic.width) / 2)
    pic.top = int(y + (h - pic.height) / 2)


def get_generated_images():
    base = ROOT / "ppt_materials" / "02_generated_samples"
    return sorted([p for p in base.rglob("*.png")])


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), DEEP_BLUE)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.18), CQU_RED)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(6.95), Inches(13.333), Inches(0.14), GOLD)
    add_text(slide, "重庆大学 SRTP 结项答辩", Inches(0.8), Inches(0.78), Inches(4.2), Inches(0.35), 15, True, RGBColor(224, 232, 245))
    add_text(slide, "基于 StyleGAN2-ADA 的\n汽车外观设计辅助系统设计与实现", Inches(0.78), Inches(1.75), Inches(8.9), Inches(1.35), 33, True, WHITE)
    add_text(slide, "三类车型可选生成 · 云端训练 · 桌面程序 · 安装包交付", Inches(0.83), Inches(3.35), Inches(7.2), Inches(0.38), 17, False, RGBColor(220, 226, 235))
    add_text(slide, "项目负责人：（请填写）\n指导教师：（请填写）\n学院专业：（请填写）\n2026 年 6 月", Inches(0.85), Inches(5.28), Inches(4.8), Inches(0.9), 14, False, WHITE)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.6), Inches(1.35), Inches(3.7), Inches(4.6), RGBColor(243, 246, 250), None, 5)
    add_text(slide, "最终成果", Inches(9.0), Inches(1.76), Inches(2.9), Inches(0.38), 20, True, CQU_RED, align=PP_ALIGN.CENTER)
    add_multiline(slide, [
        "跑车 / 轿车 / SUV 三类生成模型",
        "Win11 桌面程序可按选择生成图片",
        "Inno Setup 分卷安装包已验证",
        "源码与 Release 资产云端归档"
    ], Inches(9.05), Inches(2.42), Inches(2.82), Inches(2.2), 14, DEEP_BLUE, bullet=True)
    return slide


def slide_objectives(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide); add_header(slide, 2, "立项目标落在一个可运行系统上")
    add_card(slide, Inches(0.75), Inches(1.55), Inches(3.75), Inches(4.4), "为什么做", [
        "汽车外观概念设计早期需要大量候选图，传统草图效率依赖经验。",
        "生成式模型可提供快速视觉启发，适合辅助风格探索和方案讨论。"
    ])
    add_card(slide, Inches(4.85), Inches(1.55), Inches(3.75), Inches(4.4), "要完成什么", [
        "系统能够根据选择生成三类汽车外观图：跑车、轿车、SUV。",
        "不仅训练模型，还要做成可运行、可安装、可展示的桌面程序。"
    ], GOLD)
    add_card(slide, Inches(8.95), Inches(1.55), Inches(3.75), Inches(4.4), "最终怎么交付", [
        "三类微调模型、桌面应用、分卷安装包、GitHub 仓库、Release 资产、结项报告和答辩素材。"
    ], DEEP_BLUE)
    add_footer(slide)


def slide_route(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide); add_header(slide, 3, "技术路线从数据到安装包形成闭环")
    steps = [
        ("数据清洗", "CompCars 图像筛选、裁剪、256×256 标准化"),
        ("三类划分", "基于车型属性构建 sports / sedan / suv 数据集"),
        ("云端训练", "RTX 5090 上训练 baseline 与三类微调模型"),
        ("桌面集成", "Tkinter + PyTorch 加载对应模型生成图片"),
        ("打包归档", "PyInstaller + Inno Setup + GitHub Release")
    ]
    x0 = Inches(0.7); y = Inches(2.25); w = Inches(2.1); h = Inches(1.1); gap = Inches(0.42)
    for i, (t, b) in enumerate(steps):
        x = x0 + i * (w + gap)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h, WHITE, RGBColor(218, 225, 235))
        add_text(slide, t, x + Inches(0.12), y + Inches(0.16), w - Inches(0.24), Inches(0.25), 15, True, CQU_RED, align=PP_ALIGN.CENTER)
        add_multiline(slide, [b], x + Inches(0.18), y + Inches(0.52), w - Inches(0.36), Inches(0.45), 9.8, INK)
        if i < len(steps) - 1:
            add_arrow(slide, x + w + Inches(0.05), y + h / 2, x + w + gap - Inches(0.08), y + h / 2)
    add_metric(slide, Inches(1.0), Inches(4.45), Inches(2.2), Inches(1.15), "49,172", "清洗后通用图像", "256×256")
    add_metric(slide, Inches(3.55), Inches(4.45), Inches(2.2), Inches(1.15), "3", "车型微调模型", "sports / sedan / suv", GOLD)
    add_metric(slide, Inches(6.1), Inches(4.45), Inches(2.2), Inches(1.15), "RTX 5090", "云端训练 GPU", "32GB 显存", DEEP_BLUE)
    add_metric(slide, Inches(8.65), Inches(4.45), Inches(2.2), Inches(1.15), "Win11", "桌面安装交付", "分卷安装包", CQU_RED)
    add_footer(slide)


def slide_dataset(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide); add_header(slide, 4, "三类训练数据集支撑可选车型生成")
    data = [("跑车", "sports", 1519, CQU_RED), ("轿车", "sedan", 5492, GOLD), ("SUV", "suv", 9889, DEEP_BLUE)]
    maxv = 9889
    for i, (cn, en, val, col) in enumerate(data):
        y = Inches(1.75 + i * 1.25)
        add_text(slide, cn, Inches(0.85), y + Inches(0.05), Inches(0.7), Inches(0.3), 17, True, INK)
        add_text(slide, en, Inches(1.55), y + Inches(0.08), Inches(0.8), Inches(0.25), 10, False, MUTED)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(2.55), y, Inches(7.8), Inches(0.34), RGBColor(226, 232, 240))
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(2.55), y, Inches(7.8 * val / maxv), Inches(0.34), col)
        add_text(slide, f"{val:,} 张", Inches(10.55), y - Inches(0.02), Inches(1.2), Inches(0.3), 15, True, col)
    add_card(slide, Inches(0.85), Inches(5.32), Inches(5.5), Inches(1.0), "类别映射", [
        "2/10 -> SUV；4 -> sedan；9/11/12 -> sports；未映射类别不进入三类训练集。"
    ])
    add_card(slide, Inches(6.7), Inches(5.32), Inches(5.0), Inches(1.0), "保留策略", [
        "GitHub Release 上传实际使用的三类数据集，不上传原始 CompCars 全量数据。"
    ], DEEP_BLUE)
    add_footer(slide)


def slide_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide); add_header(slide, 5, "采用三类独立微调模型，而不是临时拼接功能")
    add_text(slide, "路线选择", Inches(0.82), Inches(1.45), Inches(2.0), Inches(0.3), 14, True, CQU_RED)
    add_multiline(slide, [
        "单一条件模型需要重构标签和训练流程，风险高、周期长。",
        "三类独立微调模型能够复用已有 baseline，工程上更稳，直接满足立项要求。"
    ], Inches(0.82), Inches(1.85), Inches(4.3), Inches(1.2), 15, INK)
    boxes = [("sports.pkl", "跑车", "低矮车身、运动姿态", CQU_RED), ("sedan.pkl", "轿车", "均衡比例、日常乘用", GOLD), ("suv.pkl", "SUV", "高车身、厚重体量", DEEP_BLUE)]
    for i, (m, cn, note, col) in enumerate(boxes):
        x = Inches(5.45 + i * 2.35)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, Inches(1.65), Inches(2.05), Inches(2.25), WHITE, RGBColor(220, 226, 235))
        add_text(slide, cn, x, Inches(1.95), Inches(2.05), Inches(0.3), 20, True, col, align=PP_ALIGN.CENTER)
        add_text(slide, m, x, Inches(2.45), Inches(2.05), Inches(0.25), 12, True, INK, align=PP_ALIGN.CENTER)
        add_multiline(slide, [note], x + Inches(0.2), Inches(2.9), Inches(1.65), Inches(0.52), 11, MUTED)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.65), Inches(11.0), Inches(1.2), RGBColor(255, 250, 240), RGBColor(231, 202, 138))
    add_text(slide, "程序逻辑：用户选择车型 -> 加载对应 pkl -> 输入 seed/truncation -> 生成并保存图片", Inches(1.25), Inches(5.02), Inches(10.5), Inches(0.32), 18, True, DEEP_BLUE, align=PP_ALIGN.CENTER)
    add_footer(slide)


def slide_cloud(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide); add_header(slide, 6, "云端 RTX 5090 解决了本地训练限制")
    add_card(slide, Inches(0.75), Inches(1.45), Inches(3.1), Inches(2.65), "云端配置", [
        "RTX 5090 32GB",
        "Ubuntu 22.04",
        "Python 3.12.3",
        "PyTorch 2.7.0+cu128",
        "CUDA Runtime 12.8"
    ])
    add_card(slide, Inches(4.15), Inches(1.45), Inches(3.1), Inches(2.65), "验证结果", [
        "cuda available = True",
        "device = NVIDIA GeForce RTX 5090",
        "capability = (12, 0)",
        "arch list 含 sm_120",
        "矩阵运算测试通过"
    ], GOLD)
    add_card(slide, Inches(7.55), Inches(1.45), Inches(4.4), Inches(2.65), "训练运行方式", [
        "项目迁移到 /root/autodl-tmp/car_srtp_project",
        "无 tmux 环境下使用 nohup 后台运行",
        "通过 nvidia-smi、pid、tail -f 日志监控"
    ], DEEP_BLUE)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(4.85), Inches(10.85), Inches(0.86), WHITE, RGBColor(220, 226, 235))
    add_text(slide, "训练闭环：冒烟测试通过 -> 第一轮训练 -> 质量评估 -> 追加训练 -> 下载最终模型与样例", Inches(1.15), Inches(5.11), Inches(10.35), Inches(0.28), 16, True, CQU_RED, align=PP_ALIGN.CENTER)
    add_footer(slide)


def slide_training(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide); add_header(slide, 7, "追加训练让三类结果达到展示要求")
    rows = [
        ("跑车", "800 kimg", "+1000 kimg", "最终 sports.pkl"),
        ("轿车", "1000 kimg", "+1500 kimg", "最终 sedan.pkl"),
        ("SUV", "1000 kimg", "+1500 kimg", "最终 suv.pkl"),
    ]
    x0 = Inches(1.0); y0 = Inches(1.75); colw = [1.3, 2.4, 2.4, 3.4]
    headers = ["车型", "第一轮训练", "追加训练", "最终产物"]
    for j, h in enumerate(headers):
        x = x0 + sum(Inches(v) for v in colw[:j])
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y0, Inches(colw[j]), Inches(0.45), DEEP_BLUE)
        add_text(slide, h, x, y0 + Inches(0.11), Inches(colw[j]), Inches(0.22), 11.5, True, WHITE, align=PP_ALIGN.CENTER)
    for i, row in enumerate(rows):
        y = y0 + Inches(0.45 + i * 0.62)
        for j, val in enumerate(row):
            x = x0 + sum(Inches(v) for v in colw[:j])
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, Inches(colw[j]), Inches(0.58), WHITE, RGBColor(218, 225, 235))
            add_text(slide, val, x + Inches(0.05), y + Inches(0.15), Inches(colw[j] - 0.1), Inches(0.22), 11.5, j == 0, INK, align=PP_ALIGN.CENTER)
    add_card(slide, Inches(1.0), Inches(4.45), Inches(5.3), Inches(1.35), "为什么追加训练", [
        "第一轮 fakes 已有车辆形状和类别差异，但最终展示质量不足；追加训练从第一轮最终模型继续，避免重复训练。"
    ])
    add_card(slide, Inches(6.65), Inches(4.45), Inches(5.0), Inches(1.35), "最终评价", [
        "每类均能筛选出 10 张以上可展示图，满足结项演示对三类生成能力的要求。"
    ], GOLD)
    add_footer(slide)


def slide_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide); add_header(slide, 8, "系统已能输出可展示汽车外观样例")
    imgs = get_generated_images()
    slots = [(Inches(0.9), Inches(1.55)), (Inches(3.95), Inches(1.55)), (Inches(7.0), Inches(1.55)), (Inches(10.05), Inches(1.55))]
    for i, (x, y) in enumerate(slots):
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, Inches(2.35), Inches(2.35), WHITE, RGBColor(218, 225, 235))
        if i < len(imgs):
            add_picture_fit(slide, imgs[i], x + Inches(0.1), y + Inches(0.1), Inches(2.15), Inches(2.15))
        else:
            add_text(slide, "后续可补充\n轿车/SUV样例", x, y + Inches(0.82), Inches(2.35), Inches(0.6), 13, True, MUTED, align=PP_ALIGN.CENTER)
        add_text(slide, f"样例 {i+1}", x, y + Inches(2.48), Inches(2.35), Inches(0.22), 10, False, MUTED, align=PP_ALIGN.CENTER)
    add_card(slide, Inches(0.95), Inches(4.75), Inches(3.7), Inches(1.25), "展示结论", [
        "三类模型输出在车身高度、轮廓比例和运动感上存在可辨别差异。"
    ])
    add_card(slide, Inches(4.85), Inches(4.75), Inches(3.7), Inches(1.25), "参数记录", [
        "程序自动记录 time、car_type、seed、truncation、image_path，便于复现实验结果。"
    ], DEEP_BLUE)
    add_card(slide, Inches(8.75), Inches(4.75), Inches(3.15), Inches(1.25), "说明", [
        "当前素材包已有跑车样例；答辩前建议再补 2-3 张轿车/SUV 截图。"
    ], GOLD)
    add_footer(slide)


def slide_app(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide); add_header(slide, 9, "桌面程序把模型变成可演示工具")
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.55), Inches(5.4), Inches(4.5), WHITE, RGBColor(218, 225, 235))
    add_text(slide, "交互界面功能", Inches(1.25), Inches(1.9), Inches(4.7), Inches(0.35), 20, True, CQU_RED)
    add_multiline(slide, [
        "车型选择：跑车 / 轿车 / SUV",
        "随机种子：输入或一键随机",
        "truncation：稳定度与多样性控制",
        "单张/批量生成",
        "自动保存图片与 metadata.csv",
        "打开输出目录、另存当前图"
    ], Inches(1.3), Inches(2.45), Inches(4.6), Inches(2.6), 15, INK, bullet=True)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.55), Inches(4.9), Inches(4.5), RGBColor(245, 247, 250), RGBColor(218, 225, 235))
    add_text(slide, "后端流程", Inches(7.35), Inches(1.9), Inches(4.2), Inches(0.35), 20, True, DEEP_BLUE)
    flow = [("选择车型", "读取 sports/sedan/suv"), ("输入参数", "seed + truncation"), ("模型推理", "PyTorch + G_ema"), ("保存结果", "PNG + CSV")]
    for i, (t, b) in enumerate(flow):
        y = Inches(2.45 + i * 0.74)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.45), y, Inches(3.95), Inches(0.48), WHITE, RGBColor(220, 226, 235))
        add_text(slide, t, Inches(7.62), y + Inches(0.12), Inches(1.2), Inches(0.2), 12, True, CQU_RED)
        add_text(slide, b, Inches(8.85), y + Inches(0.12), Inches(2.25), Inches(0.2), 11, False, INK)
    add_footer(slide)


def slide_packaging(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide); add_header(slide, 10, "打包与云端归档让项目可以继续推进")
    add_metric(slide, Inches(0.9), Inches(1.55), Inches(2.35), Inches(1.25), "3.35GB", "安装包总体积", "改为分卷上传")
    add_metric(slide, Inches(3.6), Inches(1.55), Inches(2.35), Inches(1.25), "4", "安装包文件", "exe + 3 个 bin", GOLD)
    add_metric(slide, Inches(6.3), Inches(1.55), Inches(2.35), Inches(1.25), "13", "Release 附件", "模型/数据/样例", DEEP_BLUE)
    add_metric(slide, Inches(9.0), Inches(1.55), Inches(2.35), Inches(1.25), "GitHub", "源码与资产归档", "便于迁移继续开发", CQU_RED)
    add_card(slide, Inches(0.9), Inches(3.55), Inches(5.2), Inches(1.65), "安装包", [
        "PyInstaller 打包程序目录；Inno Setup 生成 Win11 安装程序；分卷后每个文件小于 2GB，适配 GitHub Release。"
    ])
    add_card(slide, Inches(6.45), Inches(3.55), Inches(5.2), Inches(1.65), "Release 资产", [
        "上传三类模型、三类训练数据集、manifest、样例图与分卷安装包；原始数据不上传，节省云端空间。"
    ], DEEP_BLUE)
    add_footer(slide)


def slide_evaluation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide); add_header(slide, 11, "结项要求已经被系统性满足")
    rows = [
        ("能否按三类选择生成", "跑车 / 轿车 / SUV 三类模型", "达成"),
        ("是否有可运行桌面程序", "Tkinter 桌面程序与打包版均能出图", "达成"),
        ("是否能给别人安装使用", "Inno Setup 安装后应用可正常生图", "达成"),
        ("是否有生成图片样例", "每类可筛选 10 张以上展示图", "达成"),
        ("是否便于后续推进", "GitHub 源码 + Release 大资产归档", "达成"),
    ]
    x0 = Inches(0.85); y0 = Inches(1.55); widths = [3.2, 5.8, 1.55]
    headers = ["验收点", "证据", "状态"]
    for j, h in enumerate(headers):
        x = x0 + sum(Inches(v) for v in widths[:j])
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y0, Inches(widths[j]), Inches(0.42), DEEP_BLUE)
        add_text(slide, h, x, y0 + Inches(0.1), Inches(widths[j]), Inches(0.22), 11.5, True, WHITE, align=PP_ALIGN.CENTER)
    for i, row in enumerate(rows):
        y = y0 + Inches(0.42 + i * 0.58)
        for j, val in enumerate(row):
            x = x0 + sum(Inches(v) for v in widths[:j])
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, Inches(widths[j]), Inches(0.54), WHITE, RGBColor(218, 225, 235))
            color = CQU_RED if j == 2 else INK
            add_text(slide, val, x + Inches(0.06), y + Inches(0.14), Inches(widths[j] - 0.12), Inches(0.2), 10.8, j == 2, color, align=PP_ALIGN.CENTER if j == 2 else PP_ALIGN.LEFT)
    add_footer(slide)


def slide_limits(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide); add_header(slide, 12, "当前系统可展示，但仍有明确升级方向")
    add_card(slide, Inches(0.85), Inches(1.55), Inches(3.6), Inches(3.5), "当前不足", [
        "分辨率为 256×256，细节不够精细。",
        "部分随机种子仍会出现形变或局部模糊。",
        "三类独立模型便于演示，但扩展更多类别时维护成本较高。",
        "暂不支持文本、草图或局部编辑控制。"
    ], CQU_RED)
    add_card(slide, Inches(4.85), Inches(1.55), Inches(3.6), Inches(3.5), "下一步优化", [
        "构建 512×512 数据集并训练更高分辨率模型。",
        "探索单模型条件生成，减少模型数量。",
        "增加图库管理、收藏、评分、导出功能。",
        "引入 CLIP、文本提示或草图控制。"
    ], GOLD)
    add_card(slide, Inches(8.85), Inches(1.55), Inches(3.0), Inches(3.5), "项目价值", [
        "完成了从训练到应用的完整闭环。",
        "证明生成式模型可用于汽车外观早期概念辅助。",
        "为后续 SRTP 深化或毕业设计延展打下基础。"
    ], DEEP_BLUE)
    add_footer(slide)


def slide_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide); add_header(slide, 13, "总结：从模型实验走到了可交付系统")
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.55), Inches(11.2), Inches(3.2), DEEP_BLUE)
    add_text(slide, "本项目完成了汽车外观设计辅助系统的核心闭环：\n数据集构建、云端训练、三类模型、桌面程序、安装包和 GitHub 归档。", Inches(1.45), Inches(2.05), Inches(10.3), Inches(0.95), 25, True, WHITE, align=PP_ALIGN.CENTER)
    add_multiline(slide, [
        "系统已支持跑车、轿车、SUV 三类可选生成。",
        "安装版应用已验证可在 Win11 环境正常生图。",
        "项目资料已整理为报告、PPT 素材包和云端 Release 资产。"
    ], Inches(2.0), Inches(3.35), Inches(9.2), Inches(0.9), 16, RGBColor(232, 238, 247), bullet=True)
    add_text(slide, "谢谢各位老师，请批评指正", Inches(0.9), Inches(5.75), Inches(11.5), Inches(0.5), 28, True, CQU_RED, align=PP_ALIGN.CENTER)
    add_footer(slide)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_cover(prs)
    slide_objectives(prs)
    slide_route(prs)
    slide_dataset(prs)
    slide_model(prs)
    slide_cloud(prs)
    slide_training(prs)
    slide_results(prs)
    slide_app(prs)
    slide_packaging(prs)
    slide_evaluation(prs)
    slide_limits(prs)
    slide_summary(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
